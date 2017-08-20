from flask import Flask, render_template, session, redirect, url_for, escape, request
from hashlib import md5
from pptx import Presentation
import MySQLdb
from werkzeug.utils import secure_filename
from flask import send_from_directory

upload_folder = '/home/harsha/harsha/FlaskApp'



app = Flask(__name__)
app.config['upload_folder'] = upload_folder
app.secret_key = 'A0Zr98j/3yX R~XHH!jmN]LWX/,?RT'
db = MySQLdb.connect(host="localhost", user="root", passwd="harsha_sdslabs", db="pptapp")
cur = db.cursor()
app.debug = True

@app.route("/")
def main():
    if 'username' in session:
        username_session = escape(session['username']).capitalize()
        return render_template('main.html', session_user_name=username_session)
    else:
    	return render_template('index.html')

@app.route("/play", methods=['GET','POST'])
def play():

	if request.method == "GET":
		#return
		return send_from_directory('static', 'audio.mp3') 

	file = request.files['file']
	#print file
	prs = Presentation(file)
	fo = ''
	#fo = open("content.txt", "w")
	# text_runs will be populated with a list of strings,
	# one for each text run in presentation
	for slide in prs.slides:
		list_of_elements = []
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			for paragraph in shape.text_frame.paragraphs:
			        line = ''
				for run in paragraph.runs:
			              line += run.text.encode('utf-8')
                                list_of_elements.append(line)
	
                for elements in list_of_elements:
                    #print elements
                    #fo.write(elements +'\n')
                    fo += (elements + '\n')
                #fo.write('\n')
                fo += ('\n')
        #fo.close()
        #print fo
        from gtts import gTTS
        import os
        print list_of_elements
        tts = gTTS(text="".join(list_of_elements) , lang='en')
        
        filename = secure_filename("audio.mp3")
        #dump(filename)
        #tts.save(os.path.join(app.config['upload_folder', filename]))
        tts.save("/home/harsha/harsha/SmartPresenter/static/audio.mp3")
        os.system("mpg321 /home/harsha/harsha/SmartPresenter/static/audio.mp3")
        return render_template('main2.html')
	#import pyttsx
	#engine = pyttsx.init()
	#engine.setProperty('rate',0.1)
	#engine.say(list_of_elements)
	#engine.runAndWait()



@app.route("/signup",methods=['POST','GET'])
def signup():
    error = None
    if 'username' in session:
        return redirect(url_for('main'))
    elif request.method == 'POST':
        username_form = request.form['username']
        password_form = request.form['password']
        cur.execute("INSERT INTO users(name,pass) VALUES(%s , %s);",[username_form, password_form])
        session['username'] = request.form['username']
    	return render_template('main.html')
    else:
    	return render_template('index.html')   
        

@app.route("/logout")
def logout():
    session.pop('username', None)
    return redirect(url_for('signup'))  

@app.route("/login",methods=['POST','GET'])
def login():
    error = None
    if 'username' in session:
        return redirect(url_for('main'))
    if request.method == 'POST':
        username_form  = request.form['username']
        password_form  = request.form['password']
        cur.execute("SELECT COUNT(1) FROM users WHERE name = %s;", [username_form]) 
        if cur.fetchone()[0]:
            cur.execute("SELECT pass FROM users WHERE name = %s;", [username_form]) 
            for row in cur.fetchall():
                if md5(password_form).hexdigest() == row[0]:
                    session['username'] = request.form['username']
                    return redirect(url_for('main'))
                else:
                    error = "Invalid Credential"
        else:
            error = "Invalid Credential"
    return render_template('index.html', error=error)

@app.route('/images/<path:path>')
def send_js(path):
    return send_from_directory('static', path)

@app.route('/music/<path:path>')
def send_music(path):
    return send_from_directory('static', path)

if __name__ == "__main__":
    app.run()    
